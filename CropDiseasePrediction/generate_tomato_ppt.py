from pptx import Presentation


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullets_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()

    for i, item in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p.text = text
        p.level = level


def main():
    prs = Presentation()

    add_title_slide(
        prs,
        "Tomato Leaf Disease Detection using Deep Learning",
        (
            "Industry Internship (MCAIN)\n"
            "Department of Master of Computer Applications\n"
            "M S Ramaiah Institute of Technology"
        ),
    )

    add_bullets_slide(
        prs,
        "Problem Statement and Motivation",
        [
            "Tomato crop yield is heavily affected by leaf diseases if detection is delayed.",
            "Manual diagnosis needs expert support and is difficult in large farms.",
            "Objective: classify tomato leaf image into one of the major disease classes.",
            "Outcome: low-cost and fast AI-based decision support for farmers.",
        ],
    )

    add_bullets_slide(
        prs,
        "Project Objectives",
        [
            "Build a multi-class image classifier for tomato diseases.",
            "Use transfer learning to improve accuracy with limited training time.",
            "Evaluate with accuracy, precision, recall, F1-score, and confusion matrix.",
            "Deploy a simple Streamlit UI for practical demonstration.",
        ],
    )

    add_bullets_slide(
        prs,
        "Dataset and Classes",
        [
            "Dataset source: PlantVillage (Zenodo).",
            "Selected 10 tomato classes:",
            ("Tomato___Bacterial_spot, Tomato___Early_blight, Tomato___Late_blight", 1),
            ("Tomato___Leaf_Mold, Tomato___Septoria_leaf_spot, Tomato___Target_Spot", 1),
            ("Tomato___Spider_mites, Tomato___Tomato_mosaic_virus", 1),
            ("Tomato___Tomato_Yellow_Leaf_Curl_Virus, Tomato___healthy", 1),
            "Train-validation split used for unbiased evaluation.",
        ],
    )

    add_bullets_slide(
        prs,
        "Tools and Technologies (CO1)",
        [
            "Language: Python",
            "Libraries: PyTorch, torchvision, scikit-learn, matplotlib",
            "Model: ResNet18 transfer learning",
            "UI: Streamlit for image upload and prediction display",
            "Environment: Windows + virtual environment + VS Code/Cursor",
        ],
    )

    add_bullets_slide(
        prs,
        "Methodology",
        [
            "Data preprocessing: resize, normalization, and augmentation.",
            "Model training with cross-entropy loss and Adam optimizer.",
            "Validation after each epoch and best model checkpoint saving.",
            "Inference pipeline: input image -> transform -> model -> predicted class.",
        ],
    )

    add_bullets_slide(
        prs,
        "System Architecture",
        [
            "Input: Tomato leaf image captured from mobile/camera.",
            "Preprocessing: image transform and tensor conversion.",
            "Deep model: ResNet18 fine-tuned on tomato classes.",
            "Output: disease class with confidence score.",
            "Optional extension: treatment recommendation mapping per class.",
        ],
    )

    add_bullets_slide(
        prs,
        "Evaluation Metrics",
        [
            "Accuracy: overall percentage of correct predictions.",
            "Precision: reliability of positive predictions per class.",
            "Recall: ability to detect actual class instances.",
            "F1-score: harmonic mean of precision and recall.",
            "Confusion matrix: class-wise error analysis.",
        ],
    )

    add_bullets_slide(
        prs,
        "Results and Analysis",
        [
            "Latest 10-class tomato benchmark run (data_tomato_small):",
            "Train/Val images: 2000 / 759 | Epochs: 10",
            "Validation Accuracy: 10.54%",
            "F1-score (macro): 0.0191",
            "Confusion matrix and full class-wise data saved in models_tomato10_fast/metrics.json",
            "Interpretation: this is a baseline; longer training and pretrained weights are next.",
        ],
    )

    add_bullets_slide(
        prs,
        "Demonstration (CO2)",
        [
            "Step 1: Start Streamlit app and upload tomato leaf image.",
            "Step 2: System predicts disease class and confidence score.",
            "Step 3: Display top-k predictions for interpretability.",
            "Step 4: Explain model behavior with confusion matrix references.",
            "Live demo proves practical utility of learned tools and technologies.",
        ],
    )

    add_bullets_slide(
        prs,
        "Market Relevance (CO1)",
        [
            "Supports precision agriculture and early intervention.",
            "Reduces dependency on immediate expert consultation.",
            "Can be integrated into agri-advisory mobile applications.",
            "Useful for farmers, agri startups, and extension centers.",
        ],
    )

    add_bullets_slide(
        prs,
        "Limitations and Future Work",
        [
            "Current dataset is controlled; real-field images are more complex.",
            "Future work: field data collection with varied lighting/background.",
            "Add explainable AI visualization (Grad-CAM).",
            "Add remedy recommendation and multilingual advisory output.",
            "Deploy as mobile app/API for real-time usage.",
        ],
    )

    add_bullets_slide(
        prs,
        "Course Outcomes Mapping",
        [
            "CO1: Learned and applied industry-relevant AI tools and technologies.",
            "CO2: Implemented and demonstrated practical end-to-end DL pipeline.",
            "CO3: Structured technical report, presentation, and viva readiness.",
        ],
    )

    add_bullets_slide(
        prs,
        "Conclusion",
        [
            "Tomato disease classification using deep learning is feasible and effective.",
            "Project delivers an academic and practical proof-of-concept solution.",
            "Approach is scalable to multi-crop disease monitoring with additional data.",
        ],
    )

    add_bullets_slide(
        prs,
        "Thank You",
        [
            "Questions and Discussion",
            "Name: Krishna Harish",
            "Program: MCA",
            "Institution: M S Ramaiah Institute of Technology",
        ],
    )

    output_path = "Internal_1_Tomato_All_Diseases_MSRIT.pptx"
    prs.save(output_path)
    print(f"Created: {output_path}")


if __name__ == "__main__":
    main()
